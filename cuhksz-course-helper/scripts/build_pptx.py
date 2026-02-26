"""
CUHKsz Course Helper - PPTX Builder Script
Builds a standardized PPTX from structured slide content JSON.

Usage:
    python build_pptx.py <content.json> <output.pptx> [--template cs|math|stats]

Dependencies:
    pip install python-pptx
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# ── Template color/font configurations ──────────────────────────────────────

TEMPLATES = {
    "math": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x00, 0x33, 0x66),
        "body_color": RGBColor(0x00, 0x00, 0x00),
        "footer_bg": RGBColor(0x00, 0x33, 0x66),
        "footer_text": RGBColor(0xFF, 0xFF, 0xFF),
        "def_box_bg": RGBColor(0xEB, 0xF3, 0xFB),
        "def_box_border": RGBColor(0x00, 0x33, 0x66),
        "thm_box_bg": RGBColor(0xF0, 0xEE, 0xF8),
        "thm_box_border": RGBColor(0x4B, 0x00, 0x82),
        "ai_accent": RGBColor(0x00, 0x7A, 0x7A),  # teal
        "title_font": "Times New Roman",
        "body_font": "Times New Roman",
        "code_font": "Courier New",
        "school_name": "SSE, CUHK(SZ)",
    },
    "cs": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x1A, 0x2B, 0x5F),
        "body_color": RGBColor(0x2C, 0x3E, 0x50),
        "footer_bg": RGBColor(0x1A, 0x2B, 0x5F),
        "footer_text": RGBColor(0xFF, 0xFF, 0xFF),
        "def_box_bg": RGBColor(0xEB, 0xF0, 0xFF),
        "def_box_border": RGBColor(0x3B, 0x4F, 0xA3),
        "thm_box_bg": RGBColor(0xEB, 0xF0, 0xFF),
        "thm_box_border": RGBColor(0x3B, 0x4F, 0xA3),
        "ai_accent": RGBColor(0xB8, 0x86, 0x0B),  # amber
        "title_font": "Georgia",
        "body_font": "Arial",
        "code_font": "Courier New",
        "school_name": "SDS, CUHK(SZ)",
    },
    "stats": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "title_color": RGBColor(0x00, 0x33, 0x66),
        "body_color": RGBColor(0x00, 0x00, 0x00),
        "footer_bg": RGBColor(0x00, 0x33, 0x66),
        "footer_text": RGBColor(0xFF, 0xFF, 0xFF),
        "def_box_bg": RGBColor(0x2D, 0x86, 0x59),
        "def_box_border": RGBColor(0x1A, 0x5C, 0x3A),
        "def_text_color": RGBColor(0xFF, 0xFF, 0xFF),  # white text in green box
        "thm_box_bg": RGBColor(0xFF, 0xFB, 0xE6),
        "thm_box_border": RGBColor(0xD4, 0xA0, 0x17),
        "ai_accent": RGBColor(0xCC, 0x44, 0x00),  # deep orange
        "title_font": "Times New Roman",
        "body_font": "Times New Roman",
        "code_font": "Courier New",
        "school_name": "SSE, CUHK(SZ)",
    },
}

# Slide dimensions: 16:9 widescreen
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
FOOTER_H = Inches(0.35)
MARGIN_L = Inches(0.75)
MARGIN_R = Inches(0.75)
MARGIN_TOP = Inches(1.2)


def set_background(slide, color: RGBColor):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_footer(slide, cfg: dict, course_info: str, slide_num: str):
    """Add standard CUHK-SZ footer bar."""
    footer = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, SLIDE_H - FOOTER_H,
        SLIDE_W, FOOTER_H
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = cfg["footer_bg"]
    footer.line.fill.background()

    # Left: school name
    tf_left = slide.shapes.add_textbox(
        Inches(0.1), SLIDE_H - FOOTER_H,
        Inches(3), FOOTER_H
    )
    p = tf_left.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = cfg["school_name"]
    run.font.size = Pt(9)
    run.font.color.rgb = cfg["footer_text"]
    run.font.name = "Arial"

    # Center: course info
    tf_center = slide.shapes.add_textbox(
        Inches(4), SLIDE_H - FOOTER_H,
        Inches(5.33), FOOTER_H
    )
    tf_center.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf_center.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = course_info
    run.font.size = Pt(9)
    run.font.color.rgb = cfg["footer_text"]
    run.font.name = "Arial"

    # Right: slide number
    tf_right = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.5), SLIDE_H - FOOTER_H,
        Inches(1.4), FOOTER_H
    )
    tf_right.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    p = tf_right.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = slide_num
    run.font.size = Pt(9)
    run.font.color.rgb = cfg["footer_text"]
    run.font.name = "Arial"


def add_slide_title(slide, cfg: dict, title_text: str, is_ai: bool = False):
    """Add title text with navy underline rule."""
    tf = slide.shapes.add_textbox(
        MARGIN_L, Inches(0.3),
        SLIDE_W - MARGIN_L - MARGIN_R, Inches(0.8)
    )
    p = tf.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.name = cfg["title_font"]
    run.font.color.rgb = cfg["ai_accent"] if is_ai else cfg["title_color"]
    if is_ai:
        run.font.italic = True

    # Underline rule
    rule = slide.shapes.add_shape(
        1,
        MARGIN_L, Inches(1.1),
        SLIDE_W - MARGIN_L - MARGIN_R, Pt(2)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = cfg["ai_accent"] if is_ai else cfg["title_color"]
    rule.line.fill.background()


def add_body_text(slide, cfg: dict, lines: list, is_ai: bool = False,
                  top: float = None):
    """Add body text block."""
    top_pos = top if top is not None else MARGIN_TOP
    tf = slide.shapes.add_textbox(
        MARGIN_L, top_pos,
        SLIDE_W - MARGIN_L - MARGIN_R,
        SLIDE_H - top_pos - FOOTER_H - Inches(0.1)
    )
    tf.text_frame.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.text_frame.paragraphs[0]
        else:
            p = tf.text_frame.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.name = cfg["body_font"]
        run.font.color.rgb = cfg["ai_accent"] if is_ai else cfg["body_color"]
        if is_ai:
            run.font.italic = True
        p.space_after = Pt(4)


def add_box(slide, cfg: dict, label: str, content_lines: list,
            box_type: str = "definition", is_ai: bool = False):
    """Add a definition or theorem box."""
    if box_type == "definition":
        bg_color = cfg["def_box_bg"]
        border_color = cfg["def_box_border"]
        text_color = cfg.get("def_text_color", cfg["body_color"])
    else:
        bg_color = cfg["thm_box_bg"]
        border_color = cfg["thm_box_border"]
        text_color = cfg["body_color"]

    box = slide.shapes.add_shape(
        1,
        MARGIN_L, MARGIN_TOP,
        SLIDE_W - MARGIN_L - MARGIN_R,
        Inches(2.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    box.line.width = Pt(1.5)

    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.bold = True
    run.font.size = Pt(13)
    run.font.name = cfg["body_font"]
    run.font.color.rgb = cfg["ai_accent"] if is_ai else text_color

    for line in content_lines:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.name = cfg["body_font"]
        run.font.color.rgb = cfg["ai_accent"] if is_ai else text_color
        if is_ai:
            run.font.italic = True


def add_helper_tag(slide, cfg: dict):
    """Add small [Helper] tag at bottom right for AI-generated slides."""
    tf = slide.shapes.add_textbox(
        SLIDE_W - Inches(1.3),
        SLIDE_H - FOOTER_H - Inches(0.25),
        Inches(1.2), Inches(0.2)
    )
    p = tf.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "[Helper]"
    run.font.size = Pt(8)
    run.font.italic = True
    run.font.color.rgb = cfg["ai_accent"]


def build_pptx(content_json: str, output_path: str, template_name: str = "math",
               course_code: str = "", total_slides: int = 0):
    """
    Build a PPTX from structured content JSON.

    content_json: path to JSON file or JSON string
    output_path: output .pptx path
    template_name: "math", "cs", or "stats"
    course_code: e.g. "MAT3007 | Lecture 1"
    """
    cfg = TEMPLATES.get(template_name.lower(), TEMPLATES["stats"])

    # Load content
    if Path(content_json).exists():
        with open(content_json, encoding="utf-8") as f:
            data = json.load(f)
    else:
        data = json.loads(content_json)

    slides_data = data.get("slides", [])
    if not total_slides:
        total_slides = len(slides_data)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        set_background(slide, cfg["bg"])

        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        body = slide_data.get("body_text", [])
        is_ai = slide_data.get("is_ai_generated", False)
        slide_num_str = f"{i + 1} / {total_slides}"

        if slide_type == "title":
            # Large centered title slide
            tf = slide.shapes.add_textbox(
                Inches(1), Inches(2),
                SLIDE_W - Inches(2), Inches(2)
            )
            p = tf.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.name = cfg["title_font"]
            run.font.color.rgb = cfg["title_color"]

            if body:
                tf2 = slide.shapes.add_textbox(
                    Inches(1), Inches(4),
                    SLIDE_W - Inches(2), Inches(2)
                )
                tf2.text_frame.word_wrap = True
                for j, line in enumerate(body):
                    p = tf2.text_frame.paragraphs[0] if j == 0 else tf2.text_frame.add_paragraph()
                    p.alignment = PP_ALIGN.CENTER
                    run = p.add_run()
                    run.text = line
                    run.font.size = Pt(16)
                    run.font.name = cfg["body_font"]
                    run.font.color.rgb = cfg["body_color"]

        elif slide_type == "section_divider":
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = cfg["title_color"]
            tf = slide.shapes.add_textbox(
                Inches(1), Inches(2.5),
                SLIDE_W - Inches(2), Inches(2)
            )
            p = tf.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = title
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.name = cfg["title_font"]
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        elif slide_type in ("definition", "theorem", "lemma"):
            add_slide_title(slide, cfg, title, is_ai)
            label = slide_type.capitalize() + ":"
            add_box(slide, cfg, label, body, box_type=slide_type, is_ai=is_ai)
            add_footer(slide, cfg, course_code, slide_num_str)
            if is_ai:
                add_helper_tag(slide, cfg)

        else:
            # Generic content slide
            add_slide_title(slide, cfg, title, is_ai)
            add_body_text(slide, cfg, body, is_ai)
            add_footer(slide, cfg, course_code, slide_num_str)
            if is_ai:
                add_helper_tag(slide, cfg)

    prs.save(output_path)
    print(f"Saved: {output_path} ({len(slides_data)} slides, template={template_name})")
    return output_path


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python build_pptx.py <content.json> <output.pptx> [--template math|cs|stats]")
        sys.exit(1)

    content_file = sys.argv[1]
    output_file = sys.argv[2]
    template = "math"
    for i, arg in enumerate(sys.argv):
        if arg == "--template" and i + 1 < len(sys.argv):
            template = sys.argv[i + 1]

    build_pptx(content_file, output_file, template_name=template)
