"""
CUHKsz Course Helper - Content Extraction Script
Extracts structured content AND images from PPTX files.

Usage:
    python extract_content.py <input.pptx> [output.json]

Images are saved to an 'images/' subfolder next to the output JSON (or PPTX).
The JSON 'image_paths' field contains relative paths like "images/slide_01_img_01.png"
for use directly in \\includegraphics{} commands.

Output JSON structure:
{
  "source_file": "filename.pptx",
  "slide_count": N,
  "images_dir": "/absolute/path/to/images",
  "slides": [
    {
      "index": 1,
      "type": "title",
      "title": "...",
      "body_text": ["..."],
      "notes": "...",
      "has_images": true/false,
      "image_paths": ["images/slide_01_img_01.png"],
      "layout_name": "..."
    }
  ]
}
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


SLIDE_TYPE_SIGNALS = {
    "title": ["course", "lecture", "instructor", "cuhk", "university"],
    "outline": ["outline", "agenda", "contents", "today", "topics"],
    "definition": ["definition", "def.", "let ", "denote", "we say"],
    "theorem": ["theorem", "lemma", "corollary", "proposition"],
    "proof": ["proof:", "proof of", "pf."],
    "example": ["example", "ex.", "illustration"],
    "exercise": ["exercise", "problem", "homework", "hw"],
    "remark": ["remark", "note:", "observation"],
    "algorithm": ["algorithm", "pseudocode", "procedure"],
    "section_divider": [],  # detected by layout/minimal content
    "summary": ["summary", "takeaway", "conclusion", "recap"],
    "reference": ["references", "bibliography", "citation"],
}


def _get_placeholder_idx(shape) -> int:
    """Return the placeholder index, or -1 if the shape is not a placeholder."""
    try:
        pf = shape.placeholder_format
        return pf.idx if pf is not None else -1
    except (ValueError, AttributeError):
        return -1


def detect_slide_type(slide, index):
    """Heuristically detect the slide type."""
    title_text = ""
    body_text = ""

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.shape_type == 13:  # Picture — skip
            continue
        text = shape.text_frame.text.lower().strip()
        ph_idx = _get_placeholder_idx(shape)
        if ph_idx == 0:
            title_text = text
        else:
            body_text += " " + text

    combined = title_text + " " + body_text

    # First slide heuristic
    if index == 0:
        return "title"

    # Section divider: minimal content
    total_text = combined.strip()
    if len(total_text) < 30 and title_text:
        return "section_divider"

    # Signal-based detection
    for slide_type, signals in SLIDE_TYPE_SIGNALS.items():
        for signal in signals:
            if signal in combined:
                return slide_type

    return "content"  # generic fallback


def extract_text_from_shape(shape):
    """Extract text lines from a shape, preserving paragraph structure."""
    if not shape.has_text_frame:
        return []
    lines = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            lines.append(text)
    return lines


def _save_picture(shape, filepath: Path) -> bool:
    """Save a Picture shape's image blob to filepath. Returns True on success."""
    try:
        image = shape.image
        ext = image.ext.lower()
        if ext == "jpeg":
            ext = "jpg"
        # Update extension in filepath if needed
        filepath = filepath.with_suffix(f".{ext}")
        filepath.write_bytes(image.blob)
        return True
    except Exception as e:
        print(f"  Warning: could not extract image: {e}")
        return False


def extract_images_from_slide(slide, slide_index: int, images_dir: Path) -> list[str]:
    """
    Extract all Picture shapes from a slide and save them to images_dir.
    Also recurses into GroupShapes.
    Returns list of relative paths like 'images/slide_01_img_01.png'.
    """
    images_dir.mkdir(parents=True, exist_ok=True)
    image_paths = []
    img_counter = 0

    def process_shape(shape):
        nonlocal img_counter
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            img_counter += 1
            # Determine extension from image blob (may be updated in _save_picture)
            filename = f"slide_{slide_index:02d}_img_{img_counter:02d}.png"
            filepath = images_dir / filename
            if _save_picture(shape, filepath):
                # Use whichever extension _save_picture ended up writing
                actual_file = next(
                    (
                        f
                        for f in images_dir.iterdir()
                        if f.stem == filepath.stem
                    ),
                    filepath,
                )
                image_paths.append(f"images/{actual_file.name}")
        elif hasattr(shape, "shapes"):  # GroupShape — recurse
            for subshape in shape.shapes:
                process_shape(subshape)

    for shape in slide.shapes:
        process_shape(shape)

    return image_paths


def extract_pptx(input_path: str, output_path: str = None):
    prs = Presentation(input_path)
    input_path = Path(input_path).resolve()

    # images/ folder: next to output JSON, or next to the PPTX
    base_dir = Path(output_path).parent if output_path else input_path.parent
    images_dir = base_dir / "images"

    slides_data = []

    for i, slide in enumerate(prs.slides):
        slide_index = i + 1

        # Extract images for this slide
        image_paths = extract_images_from_slide(slide, slide_index, images_dir)
        has_images = bool(image_paths) or any(
            s.shape_type == 13 for s in slide.shapes
        )

        slide_info = {
            "index": slide_index,
            "type": detect_slide_type(slide, i),
            "title": "",
            "body_text": [],
            "notes": "",
            "has_images": has_images,
            "image_paths": image_paths,
            "layout_name": slide.slide_layout.name if slide.slide_layout else "",
        }

        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture — already extracted
                continue
            if shape.has_text_frame:
                ph_idx = _get_placeholder_idx(shape)
                if ph_idx == 0:
                    slide_info["title"] = shape.text_frame.text.strip()
                else:
                    slide_info["body_text"].extend(extract_text_from_shape(shape))

        # Speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                slide_info["notes"] = notes_tf.text.strip()

        slides_data.append(slide_info)

    result = {
        "source_file": str(input_path.name),
        "slide_count": len(slides_data),
        "images_dir": str(images_dir),
        "slides": slides_data,
    }

    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(result, f, ensure_ascii=False, indent=2)
        total_images = sum(len(s["image_paths"]) for s in slides_data)
        print(f"Extracted {len(slides_data)} slides -> {output_path}")
        if total_images:
            print(f"  Saved {total_images} images -> {images_dir}")
    else:
        print(json.dumps(result, ensure_ascii=False, indent=2))

    return result


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_content.py <input.pptx> [output.json]")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    extract_pptx(input_file, output_file)
